"""
Office format converters: PDF ↔ PowerPoint, PDF ↔ Excel, Excel → PDF, PowerPoint → PDF.
Uses same BaseConverter pattern as existing converters.
"""
from typing import List
from pathlib import Path
import tempfile
import os
import subprocess
import shutil

from app.utils.base_converter import BaseConverter, ConversionError

# Docker solo se importa si es necesario (cuando LibreOffice no está disponible)
docker = None


class PDFToPptxConverter(BaseConverter):
    """Convert PDF to PowerPoint: each page becomes a slide with rendered image."""

    @property
    def source_formats(self) -> List[str]:
        return ['pdf']

    @property
    def target_formats(self) -> List[str]:
        return ['pptx']

    def convert(self, input_path: str, output_path: str) -> bool:
        try:
            self.ensure_directory(output_path)
            import fitz  # PyMuPDF
            from pptx import Presentation
            from pptx.util import Inches

            doc = fitz.open(input_path)
            if len(doc) == 0:
                doc.close()
                raise ConversionError("PDF has no pages")

            prs = Presentation()
            blank_layout = prs.slide_layouts[6]  # Blank

            for page_num in range(len(doc)):
                page = doc[page_num]
                pix = page.get_pixmap(dpi=150)
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                    pix.save(tmp.name)
                    tmp_path = tmp.name
                try:
                    slide = prs.slides.add_slide(blank_layout)
                    slide.shapes.add_picture(tmp_path, Inches(0.5), Inches(0.5),
                                             width=Inches(9), height=Inches(6))
                finally:
                    try:
                        os.unlink(tmp_path)
                    except OSError:
                        pass

            doc.close()
            prs.save(output_path)
            return True
        except ImportError as e:
            raise ConversionError(
                "PDF → PowerPoint requiere PyMuPDF y python-pptx. "
                "Instala con: pip install PyMuPDF python-pptx"
            ) from e
        except Exception as e:
            raise ConversionError(f"Conversión PDF a PowerPoint falló: {str(e)}") from e


class PptxToPDFConverter(BaseConverter):
    """Convert PowerPoint to PDF using LibreOffice headless when available."""

    @property
    def source_formats(self) -> List[str]:
        return ['pptx']

    @property
    def target_formats(self) -> List[str]:
        return ['pdf']

    def convert(self, input_path: str, output_path: str) -> bool:
        self.ensure_directory(output_path)
        input_path = os.path.abspath(input_path)
        output_path = os.path.abspath(output_path)
        base_name = Path(input_path).stem

        # Prefer full path so subprocess does not depend on PATH
        libreoffice_cmd = shutil.which("libreoffice") or shutil.which("soffice")
        
        # Si LibreOffice no está instalado en el host, usar Docker local
        if not libreoffice_cmd:
            return self._convert_with_docker(input_path, output_path, base_name)
        
        # Usar LibreOffice del host si está disponible
        try:
            with tempfile.TemporaryDirectory(prefix="pptx2pdf_") as tmpdir:
                # Usar parámetros de alta calidad para PPTX→PDF (mejora calidad de texto)
                convert_to = 'pdf:impress_pdf_Export:{"UseLosslessCompression":{"type":"boolean","value":"true"},"Quality":{"type":"long","value":"100"},"ReduceImageResolution":{"type":"boolean","value":"false"}}'
                result = subprocess.run(
                    [
                        libreoffice_cmd,
                        "--headless",
                        "--convert-to", convert_to,
                        "--outdir", tmpdir,
                        input_path,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                    env={**os.environ, "HOME": tmpdir},
                )
                if result.returncode != 0:
                    raise ConversionError(
                        f"PowerPoint → PDF falló (código {result.returncode}). "
                        f"Detalle: {(result.stderr or result.stdout or '').strip()[:200]}"
                    )
                out_pdf = Path(tmpdir) / f"{base_name}.pdf"
                if not out_pdf.exists():
                    err = (result.stderr or result.stdout or "").strip()
                    if "could not be loaded" in err or "source file" in err:
                        raise ConversionError(
                            "PowerPoint → PDF requiere LibreOffice Impress. "
                            "Instala con: apt install libreoffice-impress-nogui"
                        )
                    raise ConversionError(
                        "PowerPoint → PDF no generó el archivo PDF. "
                        f"Salida: {err[:200]}"
                    )
                shutil.copy2(out_pdf, output_path)
                return True
        except subprocess.TimeoutExpired:
            raise ConversionError(
                "PowerPoint → PDF: la conversión tardó demasiado (timeout)."
            ) from None
        except ConversionError:
            raise
        except Exception as e:
            raise ConversionError(f"PowerPoint → PDF falló: {e}") from e
    
    def _convert_with_docker(self, input_path: str, output_path: str, base_name: str) -> bool:
        """Usar contenedor Docker local para convertir PPTX→PDF cuando LibreOffice no está en el host."""
        global docker
        if docker is None:
            try:
                import docker as docker_module
                docker = docker_module
            except ImportError:
                raise ConversionError(
                    "PowerPoint → PDF requiere 'docker' Python package. "
                    "Instala con: pip install docker"
                )
        
        try:
            client = docker.from_env()
        except Exception as e:
            raise ConversionError(
                "PowerPoint → PDF requiere Docker funcionando. "
                f"Error Docker: {str(e)}"
            ) from e
        
        # Buscar imagen del contenedor document-converter
        try:
            images = client.images.list(name="document-converter")
            if not images:
                # Intentar con el nombre completo de ECR
                images = client.images.list(name="766092484543.dkr.ecr.us-east-2.amazonaws.com/document-converter")
            if not images:
                raise ConversionError(
                    "No se encontró la imagen Docker 'document-converter'. "
                    "Construye con: cd ~/document-converter && ./scripts/deploy.sh"
                )
            image = images[0]
        except Exception as e:
            raise ConversionError(f"Error buscando imagen Docker: {str(e)}") from e
        
        # Montar archivos en el contenedor
        input_dir = os.path.dirname(input_path)
        output_dir = os.path.dirname(output_path)
        input_file = os.path.basename(input_path)
        
        try:
            # Ejecutar LibreOffice en el contenedor
            convert_to = 'pdf:impress_pdf_Export:{"UseLosslessCompression":{"type":"boolean","value":"true"},"Quality":{"type":"long","value":"100"},"ReduceImageResolution":{"type":"boolean","value":"false"}}'
            # containers.run() no acepta timeout; usar cliente con timeout más alto si hace falta
            client.containers.run(
                image.id,
                command=[
                    "libreoffice",
                    "--headless",
                    "--convert-to", convert_to,
                    "--outdir", "/tmp/output",
                    f"/tmp/input/{input_file}",
                ],
                volumes={
                    input_dir: {"bind": "/tmp/input", "mode": "ro"},
                    output_dir: {"bind": "/tmp/output", "mode": "rw"},
                },
                remove=True,
            )
            
            # Verificar que el archivo se generó
            expected_output = os.path.join(output_dir, f"{base_name}.pdf")
            if not os.path.exists(expected_output):
                raise ConversionError(
                    f"PowerPoint → PDF: no se generó el archivo en {expected_output}"
                )
            
            # Renombrar si es necesario
            if expected_output != output_path:
                shutil.move(expected_output, output_path)
            
            return True
        except docker.errors.ContainerError as e:
            raise ConversionError(
                f"PowerPoint → PDF falló en Docker (código {e.exit_status}). "
                f"Detalle: {str(e)[:200]}"
            ) from e
        except docker.errors.APIError as e:
            raise ConversionError(f"Error de Docker API: {str(e)}") from e
        except Exception as e:
            raise ConversionError(f"PowerPoint → PDF error con Docker: {str(e)}") from e


class PDFToExcelConverter(BaseConverter):
    """Convert PDF to Excel: extract tables and images from PDF pages into xlsx."""

    @property
    def source_formats(self) -> List[str]:
        return ['pdf']

    @property
    def target_formats(self) -> List[str]:
        return ['xlsx']

    def convert(self, input_path: str, output_path: str) -> bool:
        try:
            self.ensure_directory(output_path)
            import pdfplumber
            import fitz  # PyMuPDF
            from openpyxl import Workbook
            from openpyxl.drawing.image import Image as XLImage
            from PIL import Image as PILImage
            import io

            wb = Workbook()
            wb.remove(wb.active)

            # Tablas: primero con líneas (bordes), luego con texto (sin bordes)
            table_settings_lines = {
                "vertical_strategy": "lines",
                "horizontal_strategy": "lines",
                "snap_tolerance": 4,
                "join_tolerance": 4,
            }
            table_settings_text = {
                "vertical_strategy": "text",
                "horizontal_strategy": "text",
                "min_words_vertical": 2,
                "min_words_horizontal": 1,
            }

            with pdfplumber.open(input_path) as pdf:
                with fitz.open(input_path) as doc:
                    for page_num, page in enumerate(pdf.pages, 1):
                        # 1. Extraer tablas (probar lines primero, luego text)
                        tables = page.extract_tables(table_settings=table_settings_lines)
                        if not tables:
                            tables = page.extract_tables(table_settings=table_settings_text)
                        if not tables:
                            text = page.extract_text()
                            if text:
                                tables = [[line] for line in text.splitlines()]
                                tables = [tables] if tables else []

                        # 2. Extraer imágenes de esta página
                        fitz_page = doc[page_num - 1]
                        image_list = fitz_page.get_images(full=True)

                        # Crear hoja por página
                        title = f"Pag_{page_num}"[:31]
                        ws = wb.create_sheet(title=title)

                        current_row = 1
                        # 3. Escribir tablas
                        for tbl in tables or []:
                            if not tbl:
                                continue
                            for row in tbl:
                                clean = [str(c).strip() if c is not None else "" for c in (row or [])]
                                ws.append(clean)
                                current_row = ws.max_row
                            current_row += 1  # espacio entre tablas

                        # 4. Insertar imágenes debajo del contenido
                        max_w = 220
                        for img_idx, img in enumerate(image_list):
                            try:
                                xref = img[0]
                                base = doc.extract_image(xref)
                                img_bytes = base["image"]
                                w_img, h_img = base["width"], base["height"]
                                pil_img = PILImage.open(io.BytesIO(img_bytes))
                                if pil_img.mode in ("RGBA", "P"):
                                    pil_img = pil_img.convert("RGB")
                                if w_img > max_w:
                                    ratio = max_w / w_img
                                    new_size = (max_w, int(h_img * ratio))
                                    pil_img = pil_img.resize(new_size, PILImage.Resampling.LANCZOS)
                                buf = io.BytesIO()
                                pil_img.save(buf, format="PNG")
                                buf.seek(0)
                                xl_img = XLImage(buf)
                                xl_img.anchor = f"A{current_row}"
                                ws.add_image(xl_img)
                                current_row += max(1, int(xl_img.height / 15))
                            except Exception:
                                pass

            if not wb.sheetnames:
                ws = wb.create_sheet(title="Pag_1")
                ws.append(["Sin tablas ni texto detectado en el PDF."])
            wb.save(output_path)
            return True
        except ImportError as e:
            raise ConversionError(
                "PDF → Excel requiere pdfplumber, openpyxl, PyMuPDF y Pillow. "
                "Instala con: pip install pdfplumber openpyxl PyMuPDF pillow"
            ) from e
        except Exception as e:
            raise ConversionError(f"Conversión PDF a Excel falló: {str(e)}") from e


class ExcelToPDFConverter(BaseConverter):
    """Convert Excel to PDF: LibreOffice (colores y layout) o ReportLab como fallback."""

    @property
    def source_formats(self) -> List[str]:
        return ['xlsx', 'xls']

    @property
    def target_formats(self) -> List[str]:
        return ['pdf']

    def convert(self, input_path: str, output_path: str) -> bool:
        self.ensure_directory(output_path)
        input_path = os.path.abspath(input_path)
        output_path = os.path.abspath(output_path)
        base_name = Path(input_path).stem

        libreoffice_cmd = shutil.which("libreoffice") or shutil.which("soffice")
        if libreoffice_cmd:
            if self._convert_with_libreoffice(libreoffice_cmd, input_path, output_path, base_name):
                return True
        if self._convert_with_docker(input_path, output_path, base_name):
            return True
        ext = Path(input_path).suffix.lower()
        if ext == ".xls":
            raise ConversionError(
                "Excel (.xls) → PDF requiere LibreOffice. "
                "Instálalo en el host o usa la imagen Docker document-converter."
            )
        return self._convert_with_reportlab(input_path, output_path)

    def _convert_with_libreoffice(
        self, libreoffice_cmd: str, input_path: str, output_path: str, base_name: str
    ) -> bool:
        """LibreOffice Calc: preserva colores, formato y evita cortes de texto."""
        try:
            with tempfile.TemporaryDirectory(prefix="xlsx2pdf_") as tmpdir:
                result = subprocess.run(
                    [
                        libreoffice_cmd,
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", tmpdir,
                        input_path,
                    ],
                    capture_output=True,
                    text=True,
                    timeout=120,
                    env={**os.environ, "HOME": tmpdir},
                )
                if result.returncode != 0:
                    return False
                out_pdf = Path(tmpdir) / f"{base_name}.pdf"
                if out_pdf.exists():
                    shutil.copy2(out_pdf, output_path)
                    return True
        except (subprocess.TimeoutExpired, OSError, Exception):
            pass
        return False

    def _convert_with_docker(
        self, input_path: str, output_path: str, base_name: str
    ) -> bool:
        """Docker document-converter con LibreOffice Calc."""
        global docker
        if docker is None:
            try:
                import docker as docker_module
                docker = docker_module
            except ImportError:
                return False
        try:
            client = docker.from_env()
        except Exception:
            return False
        images = client.images.list(name="document-converter")
        if not images:
            images = client.images.list(name="766092484543.dkr.ecr.us-east-2.amazonaws.com/document-converter")
        if not images:
            return False
        image = images[0]
        input_dir = os.path.dirname(input_path)
        output_dir = os.path.dirname(output_path)
        input_file = os.path.basename(input_path)
        try:
            client.containers.run(
                image.id,
                command=[
                    "libreoffice",
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", "/tmp/output",
                    f"/tmp/input/{input_file}",
                ],
                volumes={
                    input_dir: {"bind": "/tmp/input", "mode": "ro"},
                    output_dir: {"bind": "/tmp/output", "mode": "rw"},
                },
                remove=True,
            )
            expected = os.path.join(output_dir, f"{base_name}.pdf")
            if os.path.exists(expected):
                if expected != output_path:
                    shutil.move(expected, output_path)
                return True
        except Exception:
            pass
        return False

    def _convert_with_reportlab(self, input_path: str, output_path: str) -> bool:
        """Fallback ReportLab: sin colores pero con mejor manejo de texto (wrap)."""
        try:
            self.ensure_directory(output_path)
            from openpyxl import load_workbook
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, PageBreak
            from reportlab.lib import colors
            from reportlab.lib.units import inch

            wb = load_workbook(input_path, read_only=False, data_only=True)
            doc = SimpleDocTemplate(
                output_path,
                pagesize=letter,
                rightMargin=36,
                leftMargin=36,
                topMargin=36,
                bottomMargin=36,
            )
            elements = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                col_widths = {}
                for ri, row in enumerate(ws.iter_rows()):
                    row_vals = []
                    for ci, cell in enumerate(row):
                        val = str(cell.value) if cell.value is not None else ""
                        row_vals.append(val)
                        if val:
                            col_widths[ci] = max(col_widths.get(ci, 0.8), min(2.5, 0.5 + len(val) * 0.06))
                    rows.append(row_vals)
                if not rows:
                    continue

                # Añadir colores si hay
                style_cmds = [
                    ("FONTNAME", (0, 0), (-1, -1), "Helvetica"),
                    ("FONTSIZE", (0, 0), (-1, -1), 8),
                    ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 4),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 4),
                ]
                try:
                    for ri, row in enumerate(ws.iter_rows()):
                        for ci, cell in enumerate(row):
                            if cell.fill and cell.fill.fgColor:
                                rgb = getattr(cell.fill.fgColor, "rgb", None)
                                if rgb and isinstance(rgb, str) and len(rgb) >= 6:
                                    hex_part = rgb[-6:] if len(rgb) > 6 else rgb
                                    if all(c in "0123456789ABCDEFabcdef" for c in hex_part):
                                        style_cmds.append(
                                            ("BACKGROUND", (ci, ri), (ci, ri), colors.HexColor("#" + hex_part))
                                        )
                            elif cell.fill and getattr(cell.fill, "patternType", None):
                                style_cmds.append(("BACKGROUND", (ci, ri), (ci, ri), colors.lightgrey))
                except Exception:
                    pass

                col_w = [col_widths.get(i, 1.2) * inch for i in range(max(len(r) for r in rows))]
                table = Table(rows, colWidths=col_w[:20])
                table.setStyle(TableStyle(style_cmds))
                if elements:
                    elements.append(PageBreak())
                elements.append(table)

            wb.close()
            if not elements:
                raise ConversionError("El archivo Excel no contiene datos")
            doc.build(elements)
            return True
        except ImportError as e:
            raise ConversionError(
                "Excel → PDF requiere openpyxl y reportlab. Instala: pip install openpyxl reportlab"
            ) from e
        except Exception as e:
            raise ConversionError(f"Conversión Excel a PDF falló: {str(e)}") from e
